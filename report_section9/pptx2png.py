"""Render each slide of the built pptx to PNG for visual QA.
Reads the actual shapes (geometry/fill/line/text) so the preview reflects
the real file, not a re-spec."""
import cairosvg
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_IN = 914400
SCALE = 96  # px per inch -> 1280x720

def px(emu):
    return emu / EMU_IN * SCALE

def hexof(shape, kind):
    try:
        if kind == "fill":
            if shape.fill.type == 1:
                return "#" + str(shape.fill.fore_color.rgb)
            return "none"
        else:
            if shape.line.color and shape.line.color.type is not None:
                return "#" + str(shape.line.color.rgb)
    except Exception:
        pass
    return "none"

def line_w(shape):
    try:
        return max(shape.line.width / 12700, 0.5)
    except Exception:
        return 1.0

def runs_of(shape):
    out = []
    if not shape.has_text_frame:
        return out
    for para in shape.text_frame.paragraphs:
        for r in para.runs:
            c = "#1f2937"
            try:
                if r.font.color and r.font.color.type is not None:
                    c = "#" + str(r.font.color.rgb)
            except Exception:
                pass
            out.append((r.text, (r.font.size.pt if r.font.size else 11),
                        bool(r.font.bold), c, bool(r.font.italic)))
    return out

prs = Presentation("GSCO_Section9_Diagrams.pptx")
W = px(prs.slide_width.emu if hasattr(prs.slide_width, "emu") else prs.slide_width)
H = px(prs.slide_height.emu if hasattr(prs.slide_height, "emu") else prs.slide_height)

for si, slide in enumerate(prs.slides):
    svg = [f'<svg xmlns="http://www.w3.org/2000/svg" width="{W}" height="{H}" '
           f'viewBox="0 0 {W} {H}"><rect width="{W}" height="{H}" fill="white"/>']
    # arrowhead marker
    svg.append('<defs><marker id="tri" markerWidth="9" markerHeight="9" refX="7" refY="3" '
               'orient="auto" markerUnits="userSpaceOnUse"><path d="M0,0 L7,3 L0,6 z" fill="#555555"/></marker>'
               '<marker id="trib" markerWidth="9" markerHeight="9" refX="0" refY="3" '
               'orient="auto" markerUnits="userSpaceOnUse"><path d="M7,0 L0,3 L7,6 z" fill="#555555"/></marker></defs>')
    for sh in slide.shapes:
        try:
            x, y = px(sh.left), px(sh.top)
            w, h = px(sh.width), px(sh.height)
        except Exception:
            x = y = w = h = 0
        st = sh.shape_type
        is_line = False
        try:
            if sh.auto_shape_type is not None and str(sh.auto_shape_type) == "LINE (20)":
                is_line = True
        except Exception:
            pass
        # detect line by prst
        prst = sh._element.xml
        if "<a:prstGeom prst=\"line\"" in prst:
            is_line = True
        if is_line:
            flipH = ' flipH="1"' in prst or 'flipH="1"' in prst
            flipV = 'flipV="1"' in prst
            x1, y1, x2, y2 = x, y, x + w, y + h
            if flipH:
                x1, x2 = x + w, x
            if flipV:
                y1, y2 = y + h, y
            begin = 'beginArrowType' in prst or "<a:headEnd type=\"triangle\"" in prst
            end = "<a:tailEnd type=\"triangle\"" in prst
            mk = ''
            if end:
                mk += ' marker-end="url(#tri)"'
            if begin:
                mk += ' marker-start="url(#trib)"'
            svg.append(f'<line x1="{x1:.1f}" y1="{y1:.1f}" x2="{x2:.1f}" y2="{y2:.1f}" '
                       f'stroke="#555555" stroke-width="{line_w(sh):.1f}"{mk}/>')
            continue
        # rounded rect
        if w > 0 and h > 0:
            fill = hexof(sh, "fill")
            stroke = hexof(sh, "line")
            rad = min(8, h / 4)
            svg.append(f'<rect x="{x:.1f}" y="{y:.1f}" width="{w:.1f}" height="{h:.1f}" '
                       f'rx="{rad:.1f}" fill="{fill}" stroke="{stroke}" stroke-width="{line_w(sh):.1f}"/>')
        # text
        rs = runs_of(sh)
        if rs:
            total = len(rs)
            line_h = 15
            ty = y + h / 2 - (total - 1) * line_h / 2
            for (txt, sz, bold, col, ital) in rs:
                fw = "700" if bold else "400"
                fs = "italic" if ital else "normal"
                t = (txt.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"))
                svg.append(f'<text x="{x + w/2:.1f}" y="{ty:.1f}" font-family="Arial" '
                           f'font-size="{sz*1.18:.1f}" font-weight="{fw}" font-style="{fs}" '
                           f'fill="{col}" text-anchor="middle" dominant-baseline="middle">{t}</text>')
                ty += line_h
    svg.append("</svg>")
    out = "".join(svg)
    cairosvg.svg2png(bytestring=out.encode(), write_to=f"diagram_slide{si+1}.png", output_width=int(W), output_height=int(H))
    print("wrote diagram_slide%d.png" % (si + 1))
